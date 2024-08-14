from langchain.document_loaders import PyPDFLoader
from langchain.text_splitter import RecursiveCharacterTextSplitter
from langchain.chains.summarize import load_summarize_chain
from langchain.llms import OpenAI
from pptx import Presentation
from pptx.util import Inches, Pt
from loguru import logger
import fitz  # PyMuPDF
import io
from PIL import Image
from src.paper2ppt.configs.path_config import RESOURCES_PATH

paper_path = RESOURCES_PATH / "AWQ-paper.pdf"


# Load the PDF
loader = PyPDFLoader(paper_path)
documents = loader.load()

# Split the text into chunks
text_splitter = RecursiveCharacterTextSplitter(
    separators=["\n\n", "\n", " ", ""], chunk_size=2000, chunk_overlap=200
)
texts = text_splitter.split_documents(documents)
for idx, t in enumerate(texts):
    logger.info(f"Text {idx}: {t}")


# Extract images from PDF
def extract_images(pdf_path):
    pdf_document = fitz.open(pdf_path)
    images = []
    for page_num in range(len(pdf_document)):
        page = pdf_document[page_num]
        image_list = page.get_images()
        for image_index, img in enumerate(image_list):
            xref = img[0]
            base_image = pdf_document.extract_image(xref)
            image_bytes = base_image["image"]
            image = Image.open(io.BytesIO(image_bytes))
            images.append(image)
    pdf_document.close()
    return images


# Extract images
images = extract_images(paper_path)

# Print the number of extracted images
print(f"Number of images extracted: {len(images)}")

# You can now use these images in your PowerPoint presentation
# For example, you can save them or add them to slides

# Save extracted images (optional)
for i, img in enumerate(images):
    img.save(f"extracted_image_{i}.png")

print("Images extracted successfully!")

# Rest of your code...

# Initialize the language model
llm = OpenAI(temperature=0.5)
#
# # Create a summarization chain
# chain = load_summarize_chain(llm, chain_type="map_reduce")
#
# # Summarize the paper
# summary = chain.run(texts)
#
# # Create a new PowerPoint presentation
# prs = Presentation()
#
# # Add a title slide
# title_slide_layout = prs.slide_layouts[0]
# slide = prs.slides.add_slide(title_slide_layout)
# title = slide.shapes.title
# subtitle = slide.placeholders[1]
# title.text = "Paper Summary"
# subtitle.text = "Generated with LangChain"
#
#
# # Function to add content slides
# def add_content_slide(title, content):
#     bullet_slide_layout = prs.slide_layouts[1]
#     slide = prs.slides.add_slide(bullet_slide_layout)
#     shapes = slide.shapes
#
#     title_shape = shapes.title
#     title_shape.text = title
#
#     body_shape = shapes.placeholders[1]
#     tf = body_shape.text_frame
#     tf.text = content
#
#
#     # Adjust font size if needed
#     for paragraph in tf.paragraphs:
#         paragraph.font.size = Pt(11)
#
#
# # Split the summary into sections (you may need to adjust this based on the output)
# sections = summary.split('\n\n')
#
# # Add content slides
# for i, section in enumerate(sections, 1):
#     add_content_slide(f"Summary - Part {i}", section)
#
# # Save the presentation
# prs.save('paper_summary.pptx')
#
# print("PowerPoint summary created successfully!")
